import base64
import json
import os
import uuid
import zipfile
from argparse import ArgumentParser
from concurrent.futures import ThreadPoolExecutor, as_completed

import openai
import pandas as pd
import pymupdf4llm
import requests
from Bio.PDB import PDBParser
from datasets import load_dataset
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from tenacity import retry, stop_after_attempt, wait_exponential

client = openai.OpenAI()


def handle_pptx(file_path) -> str:
    prs = Presentation(file_path)

    text = ""

    for i, slide in enumerate(prs.slides):
        text_slide = ""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_slide += run.text + "\n"
        text += f"# Slide {i}\n{text_slide}\n\n"
    return text.strip()


def handle_docx(file_path) -> str:
    doc = Document(file_path)

    text = ""

    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text.strip()


def handle_csv(file_path) -> str:
    return json.dumps(pd.read_csv(file_path).to_dict(), ensure_ascii=False)[:10000] + "..."


def handle_jsonld(file_path) -> str:
    with open(file_path, "r") as f:
        data = json.load(f)
    return json.dumps(data, ensure_ascii=False)


def handle_xlsx(file_path) -> str:
    # Function to get the color of a cell
    def get_cell_color(cell):
        fill = cell.fill
        if fill and fill.fgColor and fill.fgColor.type == "rgb":
            return fill.fgColor.rgb
        return None

    # Read the Excel file using pandas
    df = pd.read_excel(file_path)

    # Load the workbook and worksheet using openpyxl
    with open(file_path, "rb") as f:
        wb = load_workbook(f)
    ws = wb.active

    # Create a copy of the DataFrame to add comments
    result = df.copy()
    print(df)

    for row_idx, row in enumerate(ws.iter_rows(values_only=False, min_row=2, max_row=len(df) + 1)):
        for col_idx, cell in enumerate(row):
            cell_color = get_cell_color(cell)
            if cell_color and cell_color != "00000000":  # Skip default white/empty cells
                cell_value = df.iat[row_idx, col_idx]
                result.iat[row_idx, col_idx] = f"{cell_value} (Color: {cell_color})"
    return (
        "The Excel file has been transformed into the following format: each column's header (the first row) serves as a key, with the corresponding cells in that column serving as the associated values:\n"
        + json.dumps(result.to_dict(), indent=2, ensure_ascii=False)
    )


def handle_zip(file_path) -> str:
    uid = uuid.uuid4().hex
    folder = os.path.join("temp", uid)
    os.makedirs(folder, exist_ok=True)

    # Open the zip file
    with zipfile.ZipFile(file_path, "r") as zip_ref:
        # Extract all contents to the specified directory
        zip_ref.extractall(folder)
    text = ""
    for root, dirs, files in os.walk(folder):
        for f in files:
            ext = f.split(".")[-1]
            if ext in EXT_MAPPING:
                text += f"# {f}\n" + EXT_MAPPING[ext](os.path.join(root, f)) + "\n\n"
    return text


def handle_py(file_path) -> str:
    with open(file_path, "r") as f:
        return f"```python\n{f.read()}\n```"


def handle_txt(file_path) -> str:
    with open(file_path, "r") as f:
        return f.read()


@retry(stop=stop_after_attempt(3), wait=wait_exponential(multiplier=1, min=4, max=10))
def handle_image(file_path: str, question: str) -> str:
    with open(file_path, "rb") as f:
        base64_image = base64.b64encode(f.read()).decode("utf-8")
    response = client.chat.completions.create(
        messages=[
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": f'You are transforming this image to text. Here\'s a question about this image:\n"""\n{question.strip()}\n"""\n\nA team is collaborating to answer the question, but they can\'t see the image. To provide as much information as possible, you should describe this image in details (e.g., relevant text, number, element relationships, etc.). Ensure your description is complete and accurate enough for the team to do some further processing to answer the question based solely on your description! Do not try to answer the question. Only give comprehensive text description of this image.',
                    },
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"},
                    },
                ],
            },
        ],
        model="gpt-4o",
        max_tokens=4096,
        temperature=0.5,
    )
    return response.choices[0].message.content.strip()


def handle_audio(file_path: str) -> str:
    with open(file_path, "rb") as file:
        # Create a dictionary to hold the file data
        files = {"file": (file_path.split("/")[-1], file, "multipart/form-data")}

        response = requests.post("http://166.111.121.22:8848/whisper", files=files)
    return response.json()["text"]


def handle_pdf(file_path: str) -> str:
    text = pymupdf4llm.to_markdown(file_path)
    return text


def handle_pdb(file_path: str) -> str:
    parser = PDBParser()
    # Parse the structure
    structure = parser.get_structure("protein", file_path)

    # Initialize an empty string to store all information
    pdb_info = ""

    # Loop through the structure to collect information
    for model in structure:
        pdb_info += f"Model {model.id}\n"
        for chain in model:
            pdb_info += f"  Chain {chain.id}\n"
            for residue in chain:
                pdb_info += f"    Residue {residue.resname} {residue.id[1]}\n"
                for atom in residue:
                    pdb_info += f"      Atom {atom.name}, Coordinates: {atom.coord}\n"

    return pdb_info[:10000] + "..."


def handle_xml(file_path: str) -> str:
    with open(file_path, "r") as f:
        return f.read()


def process_data(data):
    result_path = f"{result_dir}/{data['task_id']}.json"

    question = data["Question"]
    additional_text = ""
    goal = question
    if data["file_name"].strip():
        ext = data["file_name"].split(".")[-1]
        func = EXT_MAPPING[ext]
        if func == handle_image:
            additional_text = func(data["file_path"], question)
        else:
            additional_text = func(data["file_path"])
        goal += f"\n\nThe user has uploaded some files. Here are the file's contents:\n\n{additional_text}"

    print(
        f"====================\nProcessing {data['task_id']}. File path: {data['file_path']}. Goal: {goal}\nLabel: {data['Final answer']}"
    )
    response = requests.post(
        "http://127.0.0.1:5050/launch_goal",
        json={
            "goal": goal,
            "team_member_names": None,
            "team_up_depth": 1,  # the depth limit of nested teaming up
            "skip_naming": False,
            "max_turns": 16,
        },
    )

    response = response.json()
    response.append(data["Final answer"])
    with open(result_path, "w") as f:
        f.write(json.dumps(response, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    EXT_MAPPING = {
        "pptx": handle_pptx,
        "csv": handle_csv,
        "jsonld": handle_jsonld,
        "xlsx": handle_xlsx,
        "xls": handle_xlsx,
        "zip": handle_zip,
        "py": handle_py,
        "jpg": handle_image,
        "docx": handle_docx,
        "mp3": handle_audio,
        "pdf": handle_pdf,
        "txt": handle_txt,
        "pdb": handle_pdb,
        "png": handle_image,
        "xml": handle_xml,
    }

    parser = ArgumentParser()
    parser.add_argument("--level", type=int, choices=[1, 2, 3])
    parser.add_argument("--max_workers", type=int, default=1)
    args = parser.parse_args()

    result_dir = f"results/gaia/level{args.level}"
    dataset = load_dataset("gaia-benchmark/GAIA", f"2023_level{args.level}")
    os.makedirs(result_dir, exist_ok=True)

    with ThreadPoolExecutor(max_workers=args.max_workers) as executor:
        futures = [executor.submit(process_data, data) for data in dataset["validation"]]
        for future in as_completed(futures):
            future.result()
